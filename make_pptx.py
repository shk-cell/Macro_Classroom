from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

GREEN       = RGBColor(0x00, 0xFF, 0x41)
GREEN_DIM   = RGBColor(0x00, 0xB3, 0x2D)
GREEN_DARK  = RGBColor(0x00, 0x3B, 0x00)
RED         = RGBColor(0xFF, 0x00, 0x40)
YELLOW      = RGBColor(0xFF, 0xD3, 0x2A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
GRAY        = RGBColor(0x33, 0x33, 0x33)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # 검정 배경
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, color=GREEN,
                align=PP_ALIGN.LEFT, font_name="Consolas"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_paragraph(tf, text, font_size=18, bold=False,
                  color=GREEN, indent=0, font_name="Consolas"):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    if indent:
        p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_content_box(slide, lines, left=Inches(0.5), top=Inches(1.8),
                    width=Inches(12.3), height=Inches(5.2)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT

        # 들여쓰기 처리
        text = line
        color = GREEN
        size = 20
        bold = False
        indent = 0

        if text.startswith("    · "):
            text = "      " + text.strip()[2:]
            color = GREEN_DIM
            size = 17
            indent = 2
        elif text.startswith("  · "):
            text = "    " + text.strip()[2:]
            color = GREEN_DIM
            size = 18
            indent = 1
        elif text.startswith("▸ ") or text.startswith("▸  "):
            pass
        elif text == "":
            pass

        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Consolas"
    return txBox


def make_header(slide, title, subtitle=None, tag=None):
    # 상단 구분선 역할의 녹색 바
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        W, Inches(1.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_DARK
    bar.line.color.rgb = GREEN
    bar.line.width = Pt(1)

    # 태그 (막 번호 등)
    if tag:
        add_textbox(slide, tag,
                    Inches(0.4), Inches(0.1),
                    Inches(4), Inches(0.5),
                    font_size=14, color=GREEN_DIM)

    # 타이틀
    add_textbox(slide, title,
                Inches(0.4), Inches(0.25),
                Inches(12), Inches(0.9),
                font_size=36, bold=True, color=GREEN,
                font_name="Consolas")

    # 서브타이틀
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(1.0),
                    Inches(12), Inches(0.45),
                    font_size=16, color=GREEN_DIM,
                    font_name="Consolas")

    # 하단 구분선
    line = slide.shapes.add_shape(
        1,
        Inches(0), Inches(1.4),
        W, Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.color.rgb = GREEN


def add_two_col(slide, left_title, left_lines, right_title, right_lines,
                top=Inches(1.6)):
    mid = Inches(6.7)
    col_w = Inches(6.0)

    # 왼쪽 타이틀
    add_textbox(slide, left_title,
                Inches(0.4), top, col_w, Inches(0.5),
                font_size=20, bold=True, color=YELLOW)

    # 오른쪽 타이틀
    add_textbox(slide, right_title,
                mid, top, col_w, Inches(0.5),
                font_size=20, bold=True, color=YELLOW)

    # 세로 구분선
    divider = slide.shapes.add_shape(
        1,
        Inches(6.55), top,
        Pt(1), Inches(5.5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = GREEN_DARK
    divider.line.color.rgb = GREEN_DARK

    # 왼쪽 내용
    lbox = slide.shapes.add_textbox(Inches(0.4), top + Inches(0.6), col_w, Inches(4.8))
    ltf = lbox.text_frame
    ltf.word_wrap = True
    first = True
    for ln in left_lines:
        p = ltf.paragraphs[0] if first else ltf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(17)
        r.font.color.rgb = GREEN if not ln.startswith("  ") else GREEN_DIM
        r.font.name = "Consolas"

    # 오른쪽 내용
    rbox = slide.shapes.add_textbox(mid, top + Inches(0.6), col_w, Inches(4.8))
    rtf = rbox.text_frame
    rtf.word_wrap = True
    first = True
    for ln in right_lines:
        p = rtf.paragraphs[0] if first else rtf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(17)
        r.font.color.rgb = GREEN if not ln.startswith("  ") else GREEN_DIM
        r.font.name = "Consolas"


# ================================================================
# 슬라이드 1 — 타이틀
# ================================================================
slide = add_slide()

add_textbox(slide, "가르치지 않아도 배운다",
            Inches(1), Inches(1.8), Inches(11.3), Inches(1.5),
            font_size=54, bold=True, color=GREEN,
            align=PP_ALIGN.CENTER, font_name="Consolas")

add_textbox(slide, "매크로 자동화 수업에서 발견한 AI 시대의 교육",
            Inches(1), Inches(3.4), Inches(11.3), Inches(0.7),
            font_size=22, color=GREEN_DIM,
            align=PP_ALIGN.CENTER, font_name="Consolas")

add_textbox(slide, "BOOKING 1000 SEATS COMPETITION",
            Inches(1), Inches(4.2), Inches(11.3), Inches(0.6),
            font_size=18, color=GREEN_DIM,
            align=PP_ALIGN.CENTER, font_name="Consolas")

# 하단 구분선
line = slide.shapes.add_shape(1, Inches(2), Inches(5.2), Inches(9.3), Pt(1))
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.color.rgb = GREEN

# ================================================================
# 슬라이드 2 — 목차
# ================================================================
slide = add_slide()
make_header(slide, "목차")

lines = [
    "01  어떤 수업이었나",
    "",
    "02  1막: 처음엔 귀여웠다",
    "",
    "03  2막: 예상치 못한 일이 벌어졌다",
    "",
    "04  3막: DB를 막자 다른 길을 찾았다",
    "",
    "05  공방의 전 과정 정리",
    "",
    "06  느낀 점",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 3 — 어떤 수업이었나
# ================================================================
slide = add_slide()
make_header(slide, "어떤 수업이었나", subtitle="수업 배경")

lines = [
    "▸  매크로 자동화를 가르치기 위해 직접 예매 사이트 제작",
    "  ·  \"1000개 좌석을 먼저 예매하는 사람이 이긴다\"",
    "▸  학생들은 AI를 활용해 자동화 프로그램을 직접 작성",
    "  ·  Python (pyautogui / Selenium) 3가지 난이도로 제공",
    "▸  선생님 vs 학생들의 공방이 시작됨",
    "",
    "처음엔 단순한 자동화 실습이었다.",
    "그런데...",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 4 — 수업 환경 소개
# ================================================================
slide = add_slide()
make_header(slide, "수업 환경 소개", subtitle="어떤 환경에서 진행됐나")

lines = [
    "▸  대상: 중학생",
    "▸  사용 도구: Python, 브라우저 개발자 콘솔, AI 도구 자유 선택",
    "▸  규칙: 방법 제한 없음, AI 사용 자유, 실패해도 재시도 가능",
    "▸  목표: 1000개 좌석을 상대보다 먼저 더 많이 예매하면 승리",
    "▸  선생님의 역할: 사이트 운영 및 방어, 직접 가르치지 않음",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 5 — 처음 예매 시스템 구조
# ================================================================
slide = add_slide()
make_header(slide, "처음 예매 시스템 구조", subtitle="Before — 초기 구조")

lines = [
    "예매 사이트 (index.html)  ──────▶  Firebase DB",
    "                                    (인증 없음)",
    "",
    "▸  웹사이트에서 Firebase에 직접 읽기 / 쓰기 가능",
    "  ·  별도 서버 없음, 인증 없음",
    "  ·  Firebase Rules 기본값: 누구나 읽기/쓰기 가능",
    "▸  이게 문제의 시작",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 6 — 1막: 처음엔 귀여웠다
# ================================================================
slide = add_slide()
make_header(slide, "1막 — 처음엔 귀여웠다", subtitle="미숙한 시도들")

add_two_col(
    slide,
    "학생 1 / 타이밍 실수",
    [
        "대기 시간을 너무 짧게 설정",
        "→ 이름 입력할 시간 없어서",
        "   예매 실패",
    ],
    "학생 2 / 좌표 실수 + 멈춤 불가",
    [
        "좌표를 잘못 잡은 채로 실행",
        "→ STOP 기능 없어서",
        "   프로그램 끄느라 우왕좌왕",
    ],
    top=Inches(1.6)
)

# 학생 3
add_textbox(slide, "학생 3 / 브라우저 폭탄",
            Inches(0.4), Inches(4.0), Inches(12), Inches(0.4),
            font_size=20, bold=True, color=YELLOW)
add_textbox(slide,
            "브라우저 창 50~100개 한번에 실행  →  과부하, 창 하나하나 직접 닫느라 혼자 고생",
            Inches(0.4), Inches(4.5), Inches(12), Inches(0.4),
            font_size=18, color=GREEN)

add_textbox(slide, '"아직 미숙하다... 고 생각했다."',
            Inches(0.4), Inches(5.2), Inches(12), Inches(0.5),
            font_size=20, color=GREEN_DIM, align=PP_ALIGN.CENTER)

# ================================================================
# 슬라이드 7 — 2막: 예상치 못한 일 (학생 4)
# ================================================================
slide = add_slide()
make_header(slide, "2막 — 예상치 못한 일이 벌어졌다", subtitle="학생 4: 시스템 구조 분석")

lines = [
    "▸  Start 버튼이 Firebase에 직접 쓰기를 한다는 것을 발견",
    "  ·  브라우저 네트워크 탭 / Firebase 콘솔 분석",
    "▸  Firebase Rules가 열려 있다는 것을 확인",
    "  ·  읽기/쓰기 권한이 누구에게나 열려 있음",
    "",
    "결과",
    "선생님이 Start를 누르기도 전에",
    "학생 혼자 먼저 게임을 시작해버림",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 8 — 기술 배경: Firebase란?
# ================================================================
slide = add_slide()
make_header(slide, "기술 배경 — Firebase란?", tag="[ 기술 배경 ]")

lines = [
    "▸  Google이 제공하는 클라우드 데이터베이스 서비스",
    "▸  코드 없이도 URL 하나로 데이터를 저장하고 읽을 수 있음",
    "▸  Rules 설정에 따라 누구나 접근 가능할 수도 있음",
    "  ·  기본값: 인증 없이 누구나 읽기/쓰기 허용",
    "",
    "이 수업에서는",
    "처음에 Rules가 열려 있었음",
    "→ 브라우저에서 누구나 직접 DB에 접근 가능한 상태",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 9 — 2막: DB에 직접 접근 (학생 5 & 6)
# ================================================================
slide = add_slide()
make_header(slide, "2막 — DB에 직접 접근하다", subtitle="학생 5 & 6 + REST API")

add_two_col(
    slide,
    "학생 5 / Firebase REST API 직접 호출",
    [
        "▸  REST API로 DB 직접 접근",
        "▸  1000개 좌석을 한 번에 동시 쓰기",
        "  ·  게임 시작 전에 이미 전부 예매",
    ],
    "학생 6 / 기존 예매 덮어쓰기 (PUT)",
    [
        "▸  남이 예매한 좌석도 자신의 이름으로 교체",
        "▸  자신 외에는 단 한 명도 예매 불가",
        "  ·  1000석 전부 독점",
    ],
    top=Inches(1.6)
)

add_textbox(slide, "REST API?",
            Inches(0.4), Inches(4.3), Inches(12), Inches(0.4),
            font_size=20, bold=True, color=YELLOW)
add_textbox(slide,
            "GET 읽기  /  POST 생성  /  PUT 덮어쓰기  /  DELETE 삭제",
            Inches(0.4), Inches(4.75), Inches(12), Inches(0.4),
            font_size=18, color=GREEN)
add_textbox(slide,
            "가능했던 이유: Firebase DB는 REST API로 접근 가능 — URL만 알면 누구나 PUT 요청으로 데이터를 덮어쓸 수 있었음",
            Inches(0.4), Inches(5.2), Inches(12), Inches(0.5),
            font_size=16, color=GREEN_DIM)

# ================================================================
# 슬라이드 10 — 선생님의 첫 번째 대응
# ================================================================
slide = add_slide()
make_header(slide, "선생님의 첫 번째 대응", subtitle="Cloudflare Worker 방어선 추가")

lines = [
    "BEFORE",
    "  예매 사이트  →  Firebase DB  (누구나 직접 접근 가능)",
    "",
    "AFTER",
    "  예매 사이트  →  Cloudflare Worker (방어선)  →  Firebase DB",
    "",
    "▸  학생은 Worker의 /claim 엔드포인트만 호출 가능",
    "▸  좌석 중복 확인 — 이미 예매된 좌석 덮어쓰기 방지",
    "▸  Firebase 직접 접근 차단",
    "",
    "그런데...",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 11 — Worker를 우회하다
# ================================================================
slide = add_slide()
make_header(slide, "Worker를 우회하다", subtitle="학생: Firebase 계정 직접 생성")

lines = [
    "▸  Cloudflare Worker가 방어선임을 파악",
    "▸  Worker를 거치지 않고 Firebase에 직접 접근할 방법 탐색",
    "▸  Firebase에 직접 계정(이메일/비밀번호)을 만들어버림",
    "  ·  해당 계정으로 인증 토큰 획득",
    "▸  토큰으로 Firebase DB에 직접 쓰기 성공",
    "  ·  → Cloudflare Worker를 완전히 우회",
    "",
    "학생 브라우저  ──── 우회! ────▶  Firebase DB (직접 접근 성공)",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 12 — 선생님의 두 번째 대응
# ================================================================
slide = add_slide()
make_header(slide, "선생님의 두 번째 대응", subtitle="Firebase Auth 추가 — 신규 가입 차단")

lines = [
    "문제: 누구나 Firebase에 계정을 만들어서 직접 접근 가능했음",
    "",
    "▸  Firebase Auth 이메일 인증 활성화",
    "▸  Firebase Rules: 인증된 계정만 쓰기 허용",
    "▸  신규 가입(계정 생성) 완전 비활성화",
    "  ·  Admin 계정 하나만 존재, 더 이상 계정 생성 불가",
    "▸  학생의 임의 계정 생성으로 우회하는 방법 차단",
    "",
    "이제 학생들이 데이터베이스에 직접 접근하는 모든 경로 차단 완료",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 13 — 3막: 멀티스레드 발견
# ================================================================
slide = add_slide()
make_header(slide, "3막 — DB를 막자 다른 길을 찾았다", subtitle="학생 6 재등장: 멀티스레드 발견")

lines = [
    "▸  DB 직접 접근 경로가 모두 막힘",
    "▸  Worker /claim 엔드포인트를 멀티스레드로 동시 다발 호출",
    "▸  본인도 모르는 사이 멀티스레드를 사용하고 있음",
    "",
    "결과",
    "하나씩 예매하던 것이 → 순식간에 수백 개 동시 예매",
    "",
    "* 따로 프로그래밍을 배운 것도 아닌데, 잘한다…",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 14 — 학생들끼리의 선의의 경쟁
# ================================================================
slide = add_slide()
make_header(slide, "학생들끼리의 선의의 경쟁", subtitle="더 빨리, 더 많이 — 다음 강자의 등장")

add_two_col(
    slide,
    "멀티스레드 학생 (기존 강자)",
    [
        "▸  Python 멀티스레드로 수백 개 동시 예매 중",
        "▸  다른 학생들보다 압도적으로 빠름",
        "▸  자신이 1등인 줄 알고 기세등등",
    ],
    "새 학생 등장 (새 강자)",
    [
        '"파이썬 안 써도 되는데요?"',
        "▸  브라우저 콘솔에서 JS로 바로 실행",
        "▸  Promise.all() 비동기 병렬 처리",
        "▸  멀티스레드보다 압도적으로 빠름",
        "▸  최종 우승!",
    ],
    top=Inches(1.6)
)

add_textbox(slide,
            "선생님을 뚫는 경쟁이 어느새 학생들끼리의 경쟁으로 발전",
            Inches(0.4), Inches(5.3), Inches(12), Inches(0.5),
            font_size=18, color=GREEN_DIM, align=PP_ALIGN.CENTER)

# ================================================================
# 슬라이드 15 — 기술 배경: Python 멀티스레드 vs Promise.all()
# ================================================================
slide = add_slide()
make_header(slide, "Python 멀티스레드 vs JS Promise.all()", subtitle="왜 브라우저 콘솔이 더 빨랐나?", tag="[ 기술 배경 ]")

add_two_col(
    slide,
    "Python 멀티스레드 (학생 6)",
    [
        "▸  여러 스레드가 번갈아 실행",
        "▸  Selenium 오버헤드 있음",
        "  ·  DOM 탐색 → 클릭 → 모달 대기",
        "  ·  → 입력 → 제출",
        "▸  요청 하나당 수십 단계 작업",
        "▸  스레드 생성 / 전환 비용 발생",
    ],
    "JS Promise.all() + fetch() (학생 7)",
    [
        "▸  JS 코드는 싱글스레드지만",
        "▸  fetch()는 브라우저 네트워킹(C++)에 위임",
        "  ·  실제 네트워크 I/O는 완전 병렬",
        "▸  DOM 조작 없음 — HTTP 패킷만 전송",
        "▸  요청 하나당 단 1단계",
        "▸  1000개가 거의 동시에 큐에 등록됨",
    ],
    top=Inches(1.6)
)

add_textbox(slide,
            "핵심: JS 싱글스레드 = JS 코드만 싱글스레드. 네트워크 I/O는 브라우저가 병렬 처리.",
            Inches(0.4), Inches(6.0), Inches(12), Inches(0.5),
            font_size=16, color=YELLOW)

# ================================================================
# 슬라이드 16 — 3막: 브라우저 콘솔 JS
# ================================================================
slide = add_slide()
make_header(slide, "3막 — 브라우저 콘솔로 더 빠르게", subtitle="학생 7: JavaScript 병렬 처리")

lines = [
    "▸  Python + Selenium보다 더 빠른 방법 탐색",
    "▸  브라우저 개발자 콘솔에서 JavaScript로 직접 실행",
    "▸  Promise.all()로 1000개 요청을 동시에 발사",
    "",
    "// 브라우저 콘솔에서 실행",
    "Promise.all(",
    "  Array.from({length: 1000}, (_, i) =>",
    "    fetch('/claim', { method: 'POST', ... })",
    "  )",
    ");",
    "",
    "결과",
    "게임 Start 누르자마자   1초도 안 돼서   1000석 전부 예매 완료",
]
add_content_box(slide, lines, top=Inches(1.6), height=Inches(5.5))

# ================================================================
# 슬라이드 17 — 선생님의 세 번째 대응과 패배
# ================================================================
slide = add_slide()
make_header(slide, "선생님의 세 번째 대응... 그리고 패배", subtitle="Rate Limit → KV 무료 한도 소진")

lines = [
    "대응: Cloudflare KV로 한 IP당 초당 100회 요청 제한 적용",
    "",
    "▸  학생들 '한꺼번에 못 하게 됐다'며 슬퍼함",
    "  ·  근데 뚫으려고 계속 시도하고 있음",
    "",
    "예상치 못한 문제",
    "▸  Cloudflare KV 무료 플랜 — 일일 쓰기 1,000회 제한",
    "▸  수정한 지 5분도 안 되어 일일 무료 한도 전부 소진",
    "▸  Rate Limit 어쩔 수 없이 제거",
    "  ·  → 학생들 다시 대량 예매 가능",
    "",
    '"선생님이 졌다."',
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 18 — 공방의 전 과정
# ================================================================
slide = add_slide()
make_header(slide, "공방의 전 과정", subtitle="시도하고 대응하고, 또 시도하고")

rounds = [
    ("1R", "Firebase Rules 열려있음 확인 → Start 없이 게임 시작",
           "Firebase Rules 수정 (인증 없으면 쓰기 불가)"),
    ("2R", "REST API로 1000석 한번에 쓰기 / 남의 좌석 덮어쓰기",
           "Cloudflare Worker 도입 + /claim 엔드포인트"),
    ("3R", "Firebase에 직접 계정 생성 → 인증 토큰으로 DB 직접 접근",
           "Firebase Auth 활성화 + 신규 가입 차단"),
    ("4R", "Worker /claim을 멀티스레드 + 브라우저 콘솔 Promise.all()",
           "Cloudflare KV Rate Limit 추가 (초당 100회)"),
]

top = Inches(1.7)
for rnd, student, teacher in rounds:
    add_textbox(slide, rnd, Inches(0.3), top, Inches(0.7), Inches(0.35),
                font_size=18, bold=True, color=YELLOW)
    add_textbox(slide, f"학생  {student}",
                Inches(1.1), top, Inches(5.5), Inches(0.35),
                font_size=15, color=GREEN)
    add_textbox(slide, f"선생님  {teacher}",
                Inches(6.8), top, Inches(6.1), Inches(0.35),
                font_size=15, color=GREEN_DIM)
    top += Inches(0.72)

add_textbox(slide, "결말  KV 무료 한도 10분 만에 소진 → Rate Limit 제거 → 학생 승리",
            Inches(0.3), top + Inches(0.1), Inches(12.5), Inches(0.4),
            font_size=17, color=RED)

# ================================================================
# 슬라이드 19 — 학생들이 스스로 발견한 기술들
# ================================================================
slide = add_slide()
make_header(slide, "학생들이 스스로 발견한 기술들", subtitle="선생님이 가르쳐준 것은 하나도 없다")

rows = [
    ("Firebase Rules 분석",      "학생 4", "브라우저 네트워크 탭·콘솔 분석"),
    ("DB REST API 직접 호출",    "학생 5", "Firebase REST API PUT 요청"),
    ("기존 레코드 덮어쓰기",     "학생 6", "인증 없이 PUT으로 기존 데이터 교체"),
    ("Firebase 계정 직접 생성",  "학생 ?", "신규 가입으로 인증 토큰 획득 후 직접 접근"),
    ("멀티스레드 병렬 요청",     "학생 6", "AI가 threading 모듈 제안"),
    ("브라우저 콘솔 JS 실행",    "학생 7", "Promise.all()로 1000개 동시"),
]

top = Inches(1.7)
for tech, who, how in rows:
    add_textbox(slide, tech,  Inches(0.3), top, Inches(3.8), Inches(0.38), font_size=16, color=GREEN)
    add_textbox(slide, who,   Inches(4.3), top, Inches(1.5), Inches(0.38), font_size=16, color=YELLOW)
    add_textbox(slide, how,   Inches(6.0), top, Inches(6.9), Inches(0.38), font_size=16, color=GREEN_DIM)
    top += Inches(0.62)

add_textbox(slide, "* 이 중 어느 것도 선생님이 먼저 가르쳐준 것이 없다",
            Inches(0.3), top + Inches(0.1), Inches(12), Inches(0.4),
            font_size=16, color=RED)

# ================================================================
# 슬라이드 20 — 이 수업이 일반 수업과 다른 점
# ================================================================
slide = add_slide()
make_header(slide, "이 수업이 일반 수업과 다른 점")

add_two_col(
    slide,
    "일반 수업",
    [
        "▸  선생님이 내용을 가르친다",
        "▸  학생은 정해진 방법으로 따라한다",
        "▸  정답이 있고, 틀리면 감점된다",
        "▸  모르면 물어보거나 포기한다",
    ],
    "이 수업",
    [
        "▸  선생님은 환경만 만든다",
        "▸  학생은 방법을 스스로 찾는다",
        "▸  정답이 없고, 실패해도 다시 시도한다",
        "▸  모르면 AI에게 물어보며 계속 나아간다",
    ],
    top=Inches(1.7)
)

add_textbox(slide, "차이를 만든 것: 이기고 싶다는 동기 하나",
            Inches(0.4), Inches(5.5), Inches(12), Inches(0.5),
            font_size=22, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ================================================================
# 슬라이드 21 — 느낀 점 1: 효능감
# ================================================================
slide = add_slide()
make_header(slide, "느낀 점 1 — 공방이 만들어낸 효능감", subtitle="막히면 더 하고 싶어지는 심리")

add_two_col(
    slide,
    "학생 시점",
    [
        "▸  선생님이 이거 막아버렸는데?",
        "▸  이거 또 뚫어보자",
        "▸  선생님이 만든 걸 뚫었다는 성취감",
        "▸  막힐수록 오히려 더 의욕이 생김",
        "▸  실패해도 포기하지 않고 다른 방법 탐색",
        "  ·  AI에게 물어보며 스스로 해결책 발견",
    ],
    "선생님 시점",
    [
        "▸  학생들이 예상을 뛰어넘는 방법을",
        "   찾아올 때마다 흥미로웠음",
        "▸  내가 만든 걸 누군가 진지하게",
        "   분석하고 있다는 것 자체가 재밌었음",
    ],
    top=Inches(1.7)
)

add_textbox(slide,
            "막고 뚫는 이 과정이 학생도, 선생님도 가장 재밌었던 순간이었다",
            Inches(0.4), Inches(5.5), Inches(12), Inches(0.5),
            font_size=18, color=GREEN_DIM, align=PP_ALIGN.CENTER)

# ================================================================
# 슬라이드 22 — 느낀 점 2: 개념은 몰라도 느낌은 안다
# ================================================================
slide = add_slide()
make_header(slide, "느낀 점 2 — 개념은 몰라도 느낌은 안다", subtitle="학생들이 이 모든 개념을 자세하게 아는 건 아니다")

convs = [
    ("대화 1 — DB 직접 접근 발견",
     '선생님: "야 너 예매 중지 눌렀는데 어떻게 뚫고 예매했어?"',
     '학생: "아, 이거 예매사이트에서 직접 데이터베이스에 접근된다고 하던데요?"',
     "→ Firebase REST API를 정확히는 모르지만, 직접 접근된다는 느낌은 알고 있음"),
    ("대화 2 — 멀티스레드",
     '선생님: "야 너 어떻게 한꺼번에 1000개씩 예매하는 거야?"',
     '학생: "저도 잘은 모르는데요, 멀티스레드 쓰면 빠르더라고요"',
     "→ 정확한 원리는 모르지만 효과는 체감함"),
]

top = Inches(1.7)
for title, q, a, note in convs:
    add_textbox(slide, title, Inches(0.4), top, Inches(12), Inches(0.35),
                font_size=17, bold=True, color=YELLOW)
    add_textbox(slide, q, Inches(0.4), top + Inches(0.4), Inches(12), Inches(0.35),
                font_size=16, color=GREEN)
    add_textbox(slide, a, Inches(0.4), top + Inches(0.78), Inches(12), Inches(0.35),
                font_size=16, color=GREEN)
    add_textbox(slide, note, Inches(0.4), top + Inches(1.16), Inches(12), Inches(0.35),
                font_size=14, color=GREEN_DIM)
    top += Inches(1.75)

add_textbox(slide, "완전히 이해하지 않아도 — AI와 함께라면 직관으로 해낸다",
            Inches(0.4), Inches(6.6), Inches(12), Inches(0.4),
            font_size=18, color=GREEN_DIM, align=PP_ALIGN.CENTER)

# ================================================================
# 슬라이드 23 — 느낀 점 2: AI의 위력
# ================================================================
slide = add_slide()
make_header(slide, "느낀 점 2 — AI의 위력", subtitle="개발 경험 유무가 무의미해졌다")

lines = [
    "▸  코딩을 전혀 모르는 학생도 AI에게 물어보며 시스템 구조를 분석",
    "▸  어떻게 하면 더 빠르게 할 수 있어? 라고 AI에 물어본 것이 전부",
    "▸  개발자 수준의 시도가 중학생에게서 나옴",
    "▸  멀티스레드, REST API, 병렬 처리",
    "  ·  배운 적 없지만 AI가 알려줬다",
    "",
    "AI를 활용하니",
    "배운 사람과 안 배운 사람의 차이가 사라졌다",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 24 — 느낀 점 3: 동기의 힘
# ================================================================
slide = add_slide()
make_header(slide, "느낀 점 3 — 동기의 힘", subtitle='"내가 해냈다"는 성취감')

lines = [
    "▸  학생들이 평소보다 훨씬 높은 집중력과 창의력을 발휘",
    "▸  실패해도 포기하지 않고 다른 방법을 스스로 탐색",
    "▸  내가 해냈다는 뿌듯함 → 강력한 학습 동기로 연결",
    "▸  서로 시도하고 대응하는 과정이 수업보다 재미있었음",
    "",
    "정답이 정해진 문제가 아니라,",
    "실제로 작동하는 시스템이어야 한다",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 25 — 느낀 점 4: AI 시대, 교사의 역할
# ================================================================
slide = add_slide()
make_header(slide, "느낀 점 4 — AI 시대, 교사의 역할", subtitle="지식 전달에서 환경 설계로")

add_two_col(
    slide,
    "이 수업에서 선생님이 직접 가르친 것",
    [
        "파이썬 매크로 프로그램",
        "기초 및 사용법만 설명",
    ],
    "학생들이 AI로 스스로 발견한 것",
    [
        "REST API",
        "멀티스레드",
        "비동기 병렬처리...",
    ],
    top=Inches(1.7)
)

lines2 = [
    "▸  AI를 활용하면 코딩 경험이 없어도 생각 이상의 성과를 낸다",
    "▸  배경 지식의 차이보다 동기의 차이가 결과를 가른다",
    "▸  교사의 역할은 가르치는 것이 아니라 하고 싶게 만드는 것이다",
]
add_content_box(slide, lines2, top=Inches(4.3), height=Inches(2.5))

# ================================================================
# 슬라이드 26 — 다른 교과·상황에 적용한다면
# ================================================================
slide = add_slide()
make_header(slide, "다른 교과·상황에 적용한다면", subtitle="환경 설계의 핵심 조건")

add_two_col(
    slide,
    "핵심 조건 3가지",
    [
        "1. 이기고 싶게 만드는 목표",
        "   경쟁, 기록 갱신, 상대 뚫기",
        "",
        "2. 실패해도 괜찮은 환경",
        "   틀려도 감점 없음, 재시도 가능",
        "",
        "3. AI를 활용할 수 있는 자유",
        "   방법 제한 없음, 도구 제한 없음",
    ],
    "교과별 적용 예시",
    [
        "▸  국어",
        "   가장 설득력 있는 글을 AI로 써서 대결",
        "",
        "▸  수학",
        "   AI로 문제를 만들고 서로 풀기",
        "",
        "▸  과학",
        "   AI로 실험 설계하고 결과 예측 대결",
    ],
    top=Inches(1.7)
)

# ================================================================
# 슬라이드 27 — 정리
# ================================================================
slide = add_slide()
make_header(slide, "정리")

lines = [
    "▸  AI를 활용하면 코딩 경험이 없어도 생각 이상의 성과를 낸다",
    "▸  배경 지식의 차이보다 동기의 차이가 결과를 가른다",
    "▸  학생 스스로 문제를 찾고 해결하게 만드는 환경이 핵심이다",
    "▸  교사의 역할은 가르치는 것이 아니라 하고 싶게 만드는 것이다",
    "",
    "AI 시대, 교사가 해야 할 일",
    "학생이 해보고 싶게 만드는 것",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 28 — 오늘 연수 안내
# ================================================================
slide = add_slide()
make_header(slide, "오늘 연수 안내", subtitle="지금까지: 사례 발표 완료")

lines = [
    "앞으로",
    "",
    "▸  2차시: 참가자 직접 학생 입장으로 전환",
    "  ·  AI 도구를 활용해 스스로 과제 해결 방법 탐색",
    "",
    "▸  3차시: 참가자 간 실전 시합",
    "  ·  실시간 순위 공개 + 상위 참가자 상품 증정 이벤트",
    "",
    "지금 체험할 사이트",
    "https://macro-classroom.shk-8b6.workers.dev/",
]
add_content_box(slide, lines)

# ================================================================
# 슬라이드 29 — Q&A
# ================================================================
slide = add_slide()

add_textbox(slide, "Q & A",
            Inches(0), Inches(2.5), W, Inches(1.5),
            font_size=72, bold=True, color=GREEN,
            align=PP_ALIGN.CENTER, font_name="Consolas")

add_textbox(slide, "질문 받겠습니다",
            Inches(0), Inches(4.2), W, Inches(0.7),
            font_size=24, color=GREEN_DIM,
            align=PP_ALIGN.CENTER, font_name="Consolas")

# ================================================================
# 슬라이드 30 — 감사합니다
# ================================================================
slide = add_slide()

add_textbox(slide, "감사합니다",
            Inches(0), Inches(2.3), W, Inches(1.5),
            font_size=64, bold=True, color=GREEN,
            align=PP_ALIGN.CENTER, font_name="Consolas")

add_textbox(slide, "가르치지 않아도 배운다",
            Inches(0), Inches(4.0), W, Inches(0.7),
            font_size=22, color=GREEN_DIM,
            align=PP_ALIGN.CENTER, font_name="Consolas")

# 하단 구분선
line = slide.shapes.add_shape(1, Inches(2), Inches(4.9), Inches(9.3), Pt(1))
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.color.rgb = GREEN

# ================================================================
# 저장
# ================================================================
out = r"C:\Users\USER\Documents\GitHub\Macro_Classroom\가르치지않아도배운다_발표_v2.pptx"
prs.save(out)
print(f"저장 완료: {out}")
print(f"슬라이드 수: {len(prs.slides)}")
