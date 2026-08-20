# -*- coding: utf-8 -*-
"""
원본 PPTX 구조를 그대로 재현하는 30슬라이드 생성기
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 ──────────────────────────────────────────────
GREEN       = RGBColor(0x00, 0xFF, 0x41)
GREEN_DIM   = RGBColor(0x00, 0xB3, 0x2D)
GREEN_DARK  = RGBColor(0x00, 0x3B, 0x00)
GREEN_BG    = RGBColor(0x05, 0x10, 0x05)
GREEN_LINE  = RGBColor(0x1A, 0x3A, 0x1A)
TITLE_BG    = RGBColor(0x00, 0x10, 0x00)
RED         = RGBColor(0xFF, 0x00, 0x40)
RED_BG      = RGBColor(0x1A, 0x00, 0x00)
YELLOW      = RGBColor(0xFF, 0xD3, 0x2A)
ORANGE      = RGBColor(0xFF, 0x9F, 0x43)
GRAY        = RGBColor(0x55, 0x55, 0x55)
BLACK       = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

I = Inches  # 단축

def new_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BLACK
    return s

def rect(slide, left, top, w, h, fill):
    shp = slide.shapes.add_shape(1, I(left), I(top), I(w), I(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp

def tb(slide, text, left, top, w, h,
       size=20, bold=False, color=GREEN,
       align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(I(left), I(top), I(w), I(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.name      = "Consolas"
    return box

def header(slide, title, subtitle=None, pagenum=None):
    """원본과 동일한 헤더 구조"""
    rect(slide, 0, 0, 13.33, 1.10, GREEN_DARK)
    rect(slide, 0, 1.10, 13.33, 0.04, GREEN)
    tb(slide, title, 0.40, 0.10, 12.50, 0.75, size=34, bold=True, color=GREEN)
    if subtitle:
        tb(slide, subtitle, 0.40, 0.72, 12.50, 0.35, size=15, color=GREEN_DIM)
    if pagenum:
        tb(slide, pagenum, 12.20, 7.10, 1.00, 0.30, size=11, color=GRAY)

def bullet(slide, text, top, color=GREEN, size=20, left=0.50):
    """▸ 메인 불릿 한 줄"""
    tb(slide, text, left, top, 12.30, 0.50, size=size, color=color)

def sub_bullet(slide, text, top, color=GREEN_DIM, size=20):
    """· 서브 불릿 한 줄"""
    tb(slide, text, 0.85, top, 11.95, 0.50, size=size, color=color)

def card(slide, left, top, w, h, accent, label, title, lines, label_color=None):
    """카드 박스 (슬라이드 4 스타일)"""
    lc = label_color or accent
    rect(slide, left, top, w, h, GREEN_BG)
    rect(slide, left, top, w, 0.06, accent)
    tb(slide, label, left, top + 0.10, w, 0.40, size=14, color=lc)
    tb(slide, title, left, top + 0.50, w, 0.60, size=19, bold=True, color=accent)
    rect(slide, left + 0.30, top + 1.15, w - 0.60, 0.02, GREEN_LINE)
    content_box = slide.shapes.add_textbox(
        I(left + 0.15), I(top + 1.30), I(w - 0.30), I(h - 1.50))
    tf = content_box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text           = line
        r.font.size      = Pt(17)
        r.font.color.rgb = GREEN_DIM
        r.font.name      = "Consolas"


# ════════════════════════════════════════════════════════════
# 슬라이드 1 — 타이틀
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 1.50, 1.80, 10.30, 4.00, TITLE_BG)
rect(s, 4.00, 4.15, 5.30, 0.03, GREEN_DIM)
tb(s, "가르치지 않아도 배운다",
   1.60, 2.00, 10.00, 1.20, size=50, bold=True, color=GREEN)
tb(s, "매크로 자동화 수업에서 발견한 AI 시대의 교육",
   1.60, 3.30, 10.00, 0.70, size=25, color=GREEN_DIM)
tb(s, "BOOKING 1000 SEATS COMPETITION",
   1.60, 4.30, 10.00, 0.50, size=15, color=GRAY)
tb(s, "1 / 30", 12.20, 7.10, 1.00, 0.30, size=11, color=GRAY)

# ════════════════════════════════════════════════════════════
# 슬라이드 2 — 목차
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "목차", pagenum="2 / 30")

items = [
    ("01", "어떤 수업이었나",           1.55),
    ("02", "1막: 처음엔 귀여웠다",      2.10),
    ("03", "2막: 예상치 못한 일이 벌어졌다", 2.65),
    ("04", "3막: DB를 막자 다른 길을 찾았다", 3.20),
    ("05", "공방의 전 과정 정리",        3.75),
    ("06", "느낀 점",                    4.30),
]
for num, text, top in items:
    tb(s, num,  0.50, top, 0.60, 0.45, size=20, bold=True, color=GREEN_DIM)
    tb(s, text, 1.20, top, 11.00, 0.45, size=20, color=GREEN)

# ════════════════════════════════════════════════════════════
# 슬라이드 3 — 어떤 수업이었나
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "어떤 수업이었나", "수업 배경", "3 / 30")

bullet(s, "▸  매크로 자동화를 가르치기 위해 직접 예매 사이트 제작", 1.50)
sub_bullet(s, '   · "1000개 좌석을 먼저 예매하는 사람이 이긴다"', 1.98)
bullet(s, "▸  학생들은 AI를 활용해 자동화 프로그램을 직접 작성", 2.46)
sub_bullet(s, "   · Python (pyautogui / Selenium) 3가지 난이도로 제공", 2.94)
bullet(s, "▸  선생님 vs 학생들의 공방이 시작됨", 6.11, color=YELLOW)

# 박스
rect(s, 0.85, 4.00, 4.50, 1.60, GREEN_DARK)
rect(s, 0.85, 4.00, 0.06, 1.60, GREEN)
tb(s, "처음엔 단순한 자동화 실습이었다.\n그런데...",
   1.05, 4.10, 4.10, 1.30, size=18, color=GREEN_DIM)

# ════════════════════════════════════════════════════════════
# 슬라이드 4 — 수업 환경 소개 (NEW)
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "수업 환경 소개", "어떤 환경에서 진행됐나", "4 / 30")

bullet(s, "▸  대상: 중학생", 1.50)
bullet(s, "▸  사용 도구: Python, 브라우저 개발자 콘솔, AI 도구 자유 선택", 1.98)
bullet(s, "▸  규칙: 방법 제한 없음, AI 사용 자유, 실패해도 재시도 가능", 2.46)
bullet(s, "▸  목표: 1000개 좌석을 상대보다 먼저 더 많이 예매하면 승리", 2.94)
bullet(s, "▸  선생님의 역할: 사이트 운영 및 방어, 직접 가르치지 않음", 3.42)

# ════════════════════════════════════════════════════════════
# 슬라이드 5 — 처음 예매 시스템 구조
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "처음 예매 시스템 구조", "Before — 초기 구조", "5 / 30")

# 왼쪽 박스 (사이트)
rect(s, 1.00, 2.40, 2.80, 1.20, GREEN_DARK)
rect(s, 1.00, 2.40, 0.08, 1.20, GREEN)
tb(s, "예매 사이트\n(index.html)", 1.15, 2.60, 2.50, 0.80, size=19, bold=True, color=GREEN)

# 화살표
tb(s, "직접 Firebase API 호출  ──────────▶", 3.90, 2.80, 5.50, 0.50, size=17, color=RED)

# 오른쪽 박스 (DB)
rect(s, 9.50, 2.40, 2.80, 1.20, RED_BG)
rect(s, 9.50, 2.40, 0.08, 1.20, RED)
tb(s, "Firebase DB\n(인증 없음)", 9.65, 2.60, 2.50, 0.80, size=19, bold=True, color=RED)

bullet(s, "▸  웹사이트에서 Firebase에 직접 읽기/쓰기 가능", 4.00, color=RED)
sub_bullet(s, "   · 별도 서버 없음 — 인증 없음", 4.48, color=RED)
sub_bullet(s, "   · Firebase Rules 기본값: 누구나 읽기/쓰기 가능", 4.96, color=RED)
bullet(s, "▸  이게 문제의 시작", 5.44, color=YELLOW)

# ════════════════════════════════════════════════════════════
# 슬라이드 6 — 1막: 처음엔 귀여웠다
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "1막 — 처음엔 귀여웠다", "미숙한 시도들", "6 / 30")

card(s, 0.40, 1.50, 4.00, 4.80, GREEN, "학생 1", "타이밍 실수",
     ["대기 시간을 너무 짧게 설정",
      "→ 이름 입력할 시간 없어서",
      "   예매 실패"])

card(s, 4.70, 1.50, 4.00, 4.80, YELLOW, "학생 2", "좌표 실수 + 멈춤 불가",
     ["좌표를 잘못 잡은 채로 실행",
      "→ STOP 기능 없어서",
      "   프로그램 끄느라 우왕좌왕"])

card(s, 9.00, 1.50, 4.00, 4.80, ORANGE, "학생 3", "브라우저 폭탄",
     ["크로미움으로 브라우저 창",
      "50개~100개씩 한번에 실행",
      "→ 과부하... 엄청난 렉",
      "→ 창 하나하나 직접 닫느라",
      "   혼자 고생하고 있음"])

tb(s, '"아직 미숙하다... 고 생각했다."', 0.50, 6.55, 12.30, 0.50, size=17, color=GRAY)

# ════════════════════════════════════════════════════════════
# 슬라이드 7 — 2막: 예상치 못한 일 (학생 4)
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "2막 — 예상치 못한 일이 벌어졌다", "학생 4: 시스템 구조 분석", "7 / 30")

bullet(s, "▸  Start 버튼이 Firebase에 직접 쓰기를 한다는 것을 발견", 1.50)
sub_bullet(s, "   · 브라우저 네트워크 탭 / Firebase 콘솔 분석", 1.98)
bullet(s, "▸  Firebase Rules가 열려 있다는 것을 확인", 2.46)
sub_bullet(s, "   · 읽기/쓰기 권한이 누구에게나 열려 있음", 2.94)

rect(s, 0.85, 4.00, 4.50, 1.80, GREEN_DARK)
rect(s, 0.85, 4.00, 0.06, 1.80, YELLOW)
tb(s, "결과\n선생님이 Start를 누르기도 전에\n학생 혼자 먼저 게임을 시작해버림",
   1.05, 4.10, 4.10, 1.60, size=18, color=YELLOW)

# ════════════════════════════════════════════════════════════
# 슬라이드 8 — 기술 배경: Firebase란? (NEW)
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "기술 배경 — Firebase란?", "[ 기술 배경 ]", "8 / 30")

bullet(s, "▸  Google이 제공하는 클라우드 데이터베이스 서비스", 1.50)
bullet(s, "▸  코드 없이도 URL 하나로 데이터를 저장하고 읽을 수 있음", 1.98)
bullet(s, "▸  Rules 설정에 따라 누구나 접근 가능할 수도 있음", 2.46)
sub_bullet(s, "   · 기본값: 인증 없이 누구나 읽기/쓰기 허용", 2.94)

rect(s, 0.85, 3.80, 5.00, 1.80, GREEN_DARK)
rect(s, 0.85, 3.80, 0.06, 1.80, RED)
tb(s, "이 수업에서는\n처음에 Rules가 열려 있었음\n→ 브라우저에서 누구나 직접 DB에 접근 가능한 상태",
   1.05, 3.90, 4.70, 1.60, size=17, color=RED)

# ════════════════════════════════════════════════════════════
# 슬라이드 9 — 2막: DB에 직접 접근 (학생 5 & 6)
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "2막 — DB에 직접 접근하다", "학생 5 & 6 + REST API", "9 / 30")

# 카드 2개
card(s, 0.40, 1.50, 5.80, 3.20, GREEN, "학생 5", "Firebase REST API 직접 호출",
     ["▸  REST API로 DB 직접 접근",
      "▸  1000개 좌석을 한 번에 동시 쓰기",
      "   · 게임 시작 전에 이미 전부 예매"])

card(s, 6.90, 1.50, 5.80, 3.20, RED, "학생 6", "기존 예매 덮어쓰기 (PUT)",
     ["▸  남이 예매한 좌석도 자신의 이름으로 교체",
      "▸  자신 외에는 단 한 명도 예매 불가",
      "   · 1000석 전부 독점"])

bullet(s, "REST API?", 5.00, color=YELLOW, size=19)
tb(s, "GET  읽기  /  POST  생성  /  PUT  덮어쓰기  /  DELETE  삭제",
   0.50, 5.48, 12.30, 0.44, size=18, color=GREEN_DIM)
tb(s, "가능했던 이유: Firebase DB는 REST API로 접근 가능 — URL만 알면 누구나 PUT 요청으로 데이터를 덮어쓸 수 있었음",
   0.50, 5.96, 12.30, 0.50, size=15, color=GRAY)

# ════════════════════════════════════════════════════════════
# 슬라이드 10 — 선생님의 첫 번째 대응
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "선생님의 첫 번째 대응", "Cloudflare Worker 방어선 추가", "10 / 30")

tb(s, "BEFORE", 0.50, 1.55, 2.00, 0.40, size=17, bold=True, color=RED)
tb(s, "예매 사이트  →  Firebase DB  (누구나 직접 접근 가능)",
   0.50, 1.95, 12.00, 0.44, size=18, color=RED)

rect(s, 0.50, 2.55, 12.30, 0.03, GREEN_LINE)

tb(s, "AFTER", 0.50, 2.70, 2.00, 0.40, size=17, bold=True, color=GREEN)
tb(s, "예매 사이트  →  Cloudflare Worker (방어선)  →  Firebase DB",
   0.50, 3.10, 12.00, 0.44, size=18, color=GREEN)

bullet(s, "▸  학생은 Worker의 /claim 엔드포인트만 호출 가능", 3.70)
bullet(s, "▸  좌석 중복 확인 — 이미 예매된 좌석 덮어쓰기 방지", 4.18)
bullet(s, "▸  Firebase 직접 접근 차단", 4.66)
bullet(s, "그런데...", 5.40, color=YELLOW)

# ════════════════════════════════════════════════════════════
# 슬라이드 11 — Worker를 우회하다
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Worker를 우회하다", "학생: Firebase 계정 직접 생성", "11 / 30")

bullet(s, "▸  Cloudflare Worker가 방어선임을 파악", 1.50)
bullet(s, "▸  Worker를 거치지 않고 Firebase에 직접 접근할 방법 탐색", 1.98)
bullet(s, "▸  Firebase에 직접 계정(이메일/비밀번호)을 만들어버림", 2.46)
sub_bullet(s, "   · 해당 계정으로 인증 토큰 획득", 2.94)
bullet(s, "▸  토큰으로 Firebase DB에 직접 쓰기 성공", 3.42)
sub_bullet(s, "   · → Cloudflare Worker를 완전히 우회", 3.90)

tb(s, "학생 브라우저  ──── 우회! ────▶  Firebase DB (직접 접근 성공)",
   0.50, 5.00, 12.30, 0.50, size=19, color=RED)

# ════════════════════════════════════════════════════════════
# 슬라이드 12 — 선생님의 두 번째 대응
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "선생님의 두 번째 대응", "Firebase Auth 추가 — 신규 가입 차단", "12 / 30")

tb(s, "문제: 누구나 Firebase에 계정을 만들어서 직접 접근 가능했음",
   0.50, 1.50, 12.30, 0.44, size=19, color=RED)

bullet(s, "▸  Firebase Auth 이메일 인증 활성화", 2.10)
bullet(s, "▸  Firebase Rules: 인증된 계정만 쓰기 허용", 2.58)
bullet(s, "▸  신규 가입(계정 생성) 완전 비활성화", 3.06)
sub_bullet(s, "   · Admin 계정 하나만 존재, 더 이상 계정 생성 불가", 3.54)
bullet(s, "▸  학생의 임의 계정 생성으로 우회하는 방법 차단", 4.02)

rect(s, 0.85, 5.00, 9.00, 0.70, GREEN_DARK)
rect(s, 0.85, 5.00, 0.06, 0.70, GREEN)
tb(s, "이제 학생들이 데이터베이스에 직접 접근하는 모든 경로 차단 완료",
   1.05, 5.10, 8.60, 0.50, size=18, color=GREEN)

# ════════════════════════════════════════════════════════════
# 슬라이드 13 — 3막: 멀티스레드 발견
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "3막 — DB를 막자 다른 길을 찾았다", "학생 6 재등장: 멀티스레드 발견", "13 / 30")

bullet(s, "▸  DB 직접 접근 경로가 모두 막힘", 1.50)
bullet(s, "▸  Worker /claim 엔드포인트를 멀티스레드로 동시 다발 호출", 1.98)
bullet(s, "▸  본인도 모르는 사이 멀티스레드를 사용하고 있음", 2.46)

rect(s, 0.85, 3.30, 7.00, 1.20, GREEN_DARK)
rect(s, 0.85, 3.30, 0.06, 1.20, YELLOW)
tb(s, "결과\n하나씩 예매하던 것이  →  순식간에 수백 개 동시 예매",
   1.05, 3.40, 6.70, 1.00, size=18, color=YELLOW)

tb(s, "* 따로 프로그래밍을 배운 것도 아닌데, 잘한다…",
   0.50, 5.20, 12.30, 0.44, size=17, color=GRAY)

# ════════════════════════════════════════════════════════════
# 슬라이드 14 — 학생들끼리의 선의의 경쟁
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "학생들끼리의 선의의 경쟁", "더 빨리, 더 많이 — 다음 강자의 등장", "14 / 30")

card(s, 0.40, 1.50, 5.80, 4.00, GREEN, "멀티스레드 학생 (기존 강자)", "",
     ["▸  Python 멀티스레드로 수백 개 동시 예매 중",
      "▸  다른 학생들보다 압도적으로 빠름",
      "▸  자신이 1등인 줄 알고 기세등등"])

card(s, 6.90, 1.50, 5.80, 4.00, YELLOW, "새 학생 등장 (새 강자)", '"파이썬 안 써도 되는데요?"',
     ["▸  브라우저 콘솔에서 JS로 바로 실행",
      "▸  Promise.all() 비동기 병렬 처리",
      "▸  멀티스레드보다 압도적으로 빠름  →  최종 우승!"])

tb(s, "선생님을 뚫는 경쟁이 어느새 학생들끼리의 경쟁으로 발전",
   0.50, 6.00, 12.30, 0.44, size=17, color=GRAY,
   align=PP_ALIGN.CENTER)
tb(s, '선생님: "솔직히 말하면... 이렇게 애들이 잘할 줄 몰랐다."',
   0.50, 6.50, 12.30, 0.44, size=17, color=GRAY,
   align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 슬라이드 15 — 기술 배경: 멀티스레드 vs Promise.all()
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Python 멀티스레드 vs JS Promise.all()", "[ 기술 배경 ] — 왜 브라우저 콘솔이 더 빨랐나?", "15 / 30")

# 왼쪽 컬럼
tb(s, "Python 멀티스레드 (학생 6)", 0.40, 1.55, 5.80, 0.44, size=19, bold=True, color=GREEN)
lines_l = [
    ("▸  여러 스레드가 번갈아 실행", 2.10, GREEN),
    ("▸  Selenium 오버헤드 있음", 2.55, GREEN),
    ("   · 요청마다: DOM 탐색 → 클릭", 3.00, GREEN_DIM),
    ("   · → 모달 대기 → 입력 → 제출", 3.40, GREEN_DIM),
    ("▸  요청 하나당 수십 단계 작업", 3.80, GREEN),
    ("▸  스레드 생성 / 전환 비용 발생", 4.25, GREEN),
]
for text, top, color in lines_l:
    tb(s, text, 0.40, top, 5.80, 0.44, size=17, color=color)

# 구분선
rect(s, 6.50, 1.50, 0.04, 5.50, GREEN_DARK)

# 오른쪽 컬럼
tb(s, "JS Promise.all() + fetch() (학생 7)", 6.70, 1.55, 6.20, 0.44, size=19, bold=True, color=YELLOW)
lines_r = [
    ("▸  JS 코드는 싱글스레드지만", 2.10, YELLOW),
    ("▸  fetch()는 브라우저 네트워킹(C++)에 위임", 2.55, YELLOW),
    ("   · 실제 네트워크 I/O는 완전 병렬", 3.00, GREEN_DIM),
    ("▸  DOM 조작 없음 — HTTP 패킷만 전송", 3.80, YELLOW),
    ("▸  요청 하나당 단 1단계", 4.25, YELLOW),
    ("▸  1000개가 거의 동시에 큐에 등록됨", 4.70, YELLOW),
]
for text, top, color in lines_r:
    tb(s, text, 6.70, top, 6.20, 0.44, size=17, color=color)

tb(s, "핵심: JS 싱글스레드 = JS 코드만 싱글스레드. 네트워크 I/O는 브라우저가 병렬 처리.",
   0.50, 6.60, 12.30, 0.44, size=16, color=YELLOW)

# ════════════════════════════════════════════════════════════
# 슬라이드 16 — 3막: 브라우저 콘솔 JS
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "3막 — 브라우저 콘솔로 더 빠르게", "학생 7: JavaScript 병렬 처리", "16 / 30")

bullet(s, "▸  Python + Selenium보다 더 빠른 방법 탐색", 1.50)
bullet(s, "▸  브라우저 개발자 콘솔에서 JavaScript로 직접 실행", 1.98)
bullet(s, "▸  Promise.all()로 1000개 요청을 동시에 발사", 2.46)

rect(s, 0.85, 3.10, 7.50, 2.20, GREEN_BG)
rect(s, 0.85, 3.10, 0.06, 2.20, GREEN_DIM)
code_lines = [
    "// 브라우저 콘솔에서 실행",
    "Promise.all(",
    "  Array.from({length: 1000}, (_, i) =>",
    "    fetch('/claim', { method: 'POST', ... })",
    "  )",
    ");",
]
for i, line in enumerate(code_lines):
    tb(s, line, 1.10, 3.20 + i * 0.33, 7.00, 0.35, size=16, color=GREEN_DIM)

rect(s, 0.85, 5.60, 9.00, 0.90, GREEN_DARK)
rect(s, 0.85, 5.60, 0.06, 0.90, YELLOW)
tb(s, "결과   게임 Start 누르자마자   1초도 안 돼서   1000석 전부 예매 완료",
   1.05, 5.75, 8.60, 0.60, size=19, color=YELLOW)

# ════════════════════════════════════════════════════════════
# 슬라이드 17 — 선생님의 세 번째 대응과 패배
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "선생님의 세 번째 대응... 그리고 패배", "Rate Limit → KV 무료 한도 소진", "17 / 30")

tb(s, "대응: Cloudflare KV로 한 IP당 초당 100회 요청 제한 적용",
   0.50, 1.55, 12.30, 0.44, size=19, color=GREEN)
bullet(s, "▸  학생들 '한꺼번에 못 하게 됐다'며 슬퍼함", 2.10)
sub_bullet(s, "   · 근데 뚫으려고 계속 시도하고 있음", 2.58)

tb(s, "예상치 못한 문제", 0.50, 3.20, 12.30, 0.44, size=19, bold=True, color=RED)
bullet(s, "▸  Cloudflare KV 무료 플랜 — 일일 쓰기 1,000회 제한", 3.68, color=RED)
bullet(s, "▸  수정한 지 5분도 안 되어 일일 무료 한도 전부 소진", 4.16, color=RED)
bullet(s, "▸  Rate Limit 어쩔 수 없이 제거", 4.64, color=RED)
sub_bullet(s, "   · → 학생들 다시 대량 예매 가능", 5.12, color=RED)

tb(s, '"선생님이 졌다."', 0.50, 6.20, 12.30, 0.44, size=22, bold=True, color=RED,
   align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 슬라이드 18 — 공방의 전 과정
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "공방의 전 과정", "시도하고 대응하고, 또 시도하고", "18 / 30")

rounds = [
    ("1R", GREEN,  "Firebase Rules 열려있음 확인 → Start 없이 게임 시작",
                   "Firebase Rules 수정 (인증 없으면 쓰기 불가)"),
    ("2R", YELLOW, "REST API로 1000석 한번에 쓰기 / 남의 좌석 덮어쓰기",
                   "Cloudflare Worker 도입 + /claim 엔드포인트"),
    ("3R", ORANGE, "Firebase에 직접 계정 생성 → 인증 토큰으로 DB 직접 접근",
                   "Firebase Auth 활성화 + 신규 가입 차단"),
    ("4R", RED,    "Worker /claim을 멀티스레드 + 브라우저 콘솔 Promise.all()",
                   "Cloudflare KV Rate Limit 추가 (초당 100회)"),
]

top = 1.50
for rnd, color, student, teacher in rounds:
    tb(s, rnd,     0.30, top, 0.60, 0.38, size=18, bold=True, color=color)
    tb(s, f"학생  {student}", 1.00, top, 6.10, 0.38, size=15, color=GREEN)
    tb(s, f"선생님  {teacher}", 7.20, top, 5.80, 0.38, size=15, color=GREEN_DIM)
    rect(s, 0.30, top + 0.44, 12.70, 0.02, GREEN_LINE)
    top += 0.72 + 0.10

rect(s, 0.30, top + 0.10, 12.70, 0.60, RED_BG)
rect(s, 0.30, top + 0.10, 0.06, 0.60, RED)
tb(s, "결말  KV 무료 한도 10분 만에 소진  →  Rate Limit 제거  →  학생 승리",
   0.50, top + 0.20, 12.30, 0.40, size=17, color=RED)

# ════════════════════════════════════════════════════════════
# 슬라이드 19 — 학생들이 스스로 발견한 기술들
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "학생들이 스스로 발견한 기술들", "선생님이 가르쳐준 것은 하나도 없다", "19 / 30")

rows = [
    ("Firebase Rules 분석",     "학생 4", "브라우저 네트워크 탭·콘솔 분석"),
    ("DB REST API 직접 호출",   "학생 5", "Firebase REST API PUT 요청"),
    ("기존 레코드 덮어쓰기",    "학생 6", "인증 없이 PUT으로 기존 데이터 교체"),
    ("Firebase 계정 직접 생성", "학생 ?", "신규 가입으로 인증 토큰 획득 후 직접 접근"),
    ("멀티스레드 병렬 요청",    "학생 6", "AI가 threading 모듈 제안"),
    ("브라우저 콘솔 JS 실행",   "학생 7", "Promise.all()로 1000개 동시"),
]
# 헤더행
tb(s, "기술",   0.30, 1.52, 3.80, 0.38, size=15, bold=True, color=GRAY)
tb(s, "학생",   4.30, 1.52, 1.50, 0.38, size=15, bold=True, color=GRAY)
tb(s, "발견 방법", 6.00, 1.52, 6.90, 0.38, size=15, bold=True, color=GRAY)
rect(s, 0.30, 1.93, 12.70, 0.02, GREEN_LINE)

top = 2.00
for tech, who, how in rows:
    tb(s, tech, 0.30, top, 3.80, 0.40, size=17, color=GREEN)
    tb(s, who,  4.30, top, 1.50, 0.40, size=17, color=YELLOW)
    tb(s, how,  6.00, top, 6.90, 0.40, size=17, color=GREEN_DIM)
    top += 0.62

tb(s, "* 이 중 어느 것도 선생님이 먼저 가르쳐준 것이 없다",
   0.30, top + 0.10, 12.00, 0.40, size=16, color=RED)

# ════════════════════════════════════════════════════════════
# 슬라이드 20 — 이 수업이 일반 수업과 다른 점 (NEW)
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "이 수업이 일반 수업과 다른 점", pagenum="20 / 30")

card(s, 0.40, 1.50, 5.80, 3.80, GREEN_DIM, "", "일반 수업",
     ["▸  선생님이 내용을 가르친다",
      "▸  학생은 정해진 방법으로 따라한다",
      "▸  정답이 있고, 틀리면 감점된다",
      "▸  모르면 물어보거나 포기한다"])

card(s, 6.90, 1.50, 5.80, 3.80, GREEN, "", "이 수업",
     ["▸  선생님은 환경만 만든다",
      "▸  학생은 방법을 스스로 찾는다",
      "▸  정답이 없고, 실패해도 다시 시도한다",
      "▸  모르면 AI에게 물어보며 계속 나아간다"])

tb(s, "차이를 만든 것: 이기고 싶다는 동기 하나",
   0.50, 5.70, 12.30, 0.60, size=24, bold=True, color=YELLOW,
   align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 슬라이드 21 — 느낀 점 1: 효능감
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "느낀 점 1 — 공방이 만들어낸 효능감", "막히면 더 하고 싶어지는 심리", "21 / 30")

# 왼쪽
tb(s, "학생 시점", 0.40, 1.55, 5.80, 0.44, size=19, bold=True, color=GREEN)
ls = [
    "▸  선생님이 이거 막아버렸는데?",
    "▸  이거 또 뚫어보자",
    "▸  선생님이 만든 걸 뚫었다는 성취감",
    "▸  막힐수록 오히려 더 의욕이 생김",
    "▸  실패해도 포기하지 않고 다른 방법 탐색",
    "   · AI에게 물어보며 스스로 해결책 발견",
]
for i, line in enumerate(ls):
    color = GREEN_DIM if line.startswith("   ·") else GREEN
    tb(s, line, 0.40, 2.10 + i * 0.48, 5.80, 0.44, size=17, color=color)

# 구분선
rect(s, 6.50, 1.50, 0.04, 5.00, GREEN_LINE)

# 오른쪽
tb(s, "선생님 시점", 6.70, 1.55, 6.20, 0.44, size=19, bold=True, color=YELLOW)
rs = [
    "▸  학생들이 예상을 뛰어넘는 방법을",
    "   찾아올 때마다 흥미로웠음",
    "▸  내가 만든 걸 누군가 진지하게",
    "   분석하고 있다는 것 자체가 재밌었음",
]
for i, line in enumerate(rs):
    color = GREEN_DIM if line.startswith("   ") else YELLOW
    tb(s, line, 6.70, 2.10 + i * 0.48, 6.20, 0.44, size=17, color=color)

tb(s, "막고 뚫는 이 과정이 학생도, 선생님도 가장 재밌었던 순간이었다",
   0.50, 6.55, 12.30, 0.44, size=17, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 슬라이드 22 — 느낀 점 2: 개념은 몰라도 느낌은 안다
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "느낀 점 2 — 개념은 몰라도 느낌은 안다",
       "학생들이 이 모든 개념을 자세하게 아는 건 아니다", "22 / 30")

convs = [
    ("대화 1 — DB 직접 접근 발견", 1.50,
     '선생님: "야 너 예매 중지 눌렀는데 어떻게 뚫고 예매했어?"',
     '학생: "아, 이거 예매사이트에서 직접 데이터베이스에 접근된다고 하던데요?"',
     "→ Firebase REST API를 정확히는 모르지만, 직접 접근된다는 느낌은 알고 있음"),
    ("대화 2 — 멀티스레드", 3.70,
     '선생님: "야 너 어떻게 한꺼번에 1000개씩 예매하는 거야?"',
     '학생: "저도 잘은 모르는데요, 멀티스레드 쓰면 빠르더라고요"',
     "→ 정확한 원리는 모르지만 효과는 체감함"),
]
for title, top, q, a, note in convs:
    tb(s, title, 0.40, top, 12.30, 0.38, size=16, bold=True, color=YELLOW)
    tb(s, q, 0.40, top + 0.44, 12.30, 0.38, size=16, color=GREEN)
    tb(s, a, 0.40, top + 0.84, 12.30, 0.38, size=16, color=GREEN)
    tb(s, note, 0.40, top + 1.24, 12.30, 0.38, size=14, color=GREEN_DIM)

tb(s, "완전히 이해하지 않아도 — AI와 함께라면 직관으로 해낸다",
   0.50, 6.70, 12.30, 0.40, size=17, color=GREEN_DIM, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 슬라이드 23 — 느낀 점 2: AI의 위력
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "느낀 점 2 — AI의 위력", "개발 경험 유무가 무의미해졌다", "23 / 30")

bullet(s, "▸  코딩을 전혀 모르는 학생도 AI에게 물어보며 시스템 구조를 분석", 1.50)
bullet(s, "▸  어떻게 하면 더 빠르게 할 수 있어? 라고 AI에 물어본 것이 전부", 1.98)
bullet(s, "▸  개발자 수준의 시도가 중학생에게서 나옴", 2.46)
bullet(s, "▸  멀티스레드, REST API, 병렬 처리", 2.94)
sub_bullet(s, "   · 배운 적 없지만 AI가 알려줬다", 3.42)

rect(s, 0.85, 4.30, 9.00, 1.40, GREEN_DARK)
rect(s, 0.85, 4.30, 0.06, 1.40, GREEN)
tb(s, "AI를 활용하니\n배운 사람과 안 배운 사람의 차이가 사라졌다",
   1.05, 4.45, 8.60, 1.10, size=22, bold=True, color=GREEN)

# ════════════════════════════════════════════════════════════
# 슬라이드 24 — 느낀 점 3: 동기의 힘
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "느낀 점 3 — 동기의 힘", '"내가 해냈다"는 성취감', "24 / 30")

bullet(s, "▸  학생들이 평소보다 훨씬 높은 집중력과 창의력을 발휘", 1.50)
bullet(s, "▸  실패해도 포기하지 않고 다른 방법을 스스로 탐색", 1.98)
bullet(s, '▸  "내가 해냈다"는 뿌듯함  →  강력한 학습 동기로 연결', 2.46)
bullet(s, "▸  서로 시도하고 대응하는 과정이 수업보다 재미있었음", 2.94)

rect(s, 0.85, 4.10, 9.00, 1.40, GREEN_DARK)
rect(s, 0.85, 4.10, 0.06, 1.40, YELLOW)
tb(s, "정답이 정해진 문제가 아니라,\n실제로 작동하는 시스템이어야 한다",
   1.05, 4.25, 8.60, 1.10, size=22, bold=True, color=YELLOW)

# ════════════════════════════════════════════════════════════
# 슬라이드 25 — 느낀 점 4: AI 시대, 교사의 역할
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "느낀 점 4 — AI 시대, 교사의 역할", "지식 전달에서 환경 설계로", "25 / 30")

bullet(s, "▸  AI를 활용하면 코딩 경험이 없어도 생각 이상의 성과를 낸다", 1.50)
bullet(s, "▸  배경 지식의 차이보다 동기의 차이가 결과를 가른다", 1.98)
bullet(s, "▸  학생 스스로 문제를 찾고 해결하게 만드는 환경이 핵심이다", 2.46)
bullet(s, "▸  교사의 역할은 가르치는 것이 아니라 하고 싶게 만드는 것이다", 2.94)

# 비교 박스 2개
card(s, 0.40, 3.80, 5.50, 2.60, GREEN_DIM, "", "이 수업에서 선생님이 직접 가르친 것",
     ["파이썬 매크로 프로그램", "기초 및 사용법만 설명"])

card(s, 7.10, 3.80, 5.50, 2.60, GREEN, "", "학생들이 AI로 스스로 발견한 것",
     ["REST API", "멀티스레드", "비동기 병렬처리..."])

tb(s, "vs", 6.10, 4.70, 0.80, 0.60, size=28, bold=True, color=GRAY,
   align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 슬라이드 26 — 다른 교과에 적용한다면 (NEW)
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "다른 교과·상황에 적용한다면", "환경 설계의 핵심 조건", "26 / 30")

# 왼쪽: 핵심 조건
tb(s, "핵심 조건 3가지", 0.40, 1.55, 5.80, 0.44, size=19, bold=True, color=YELLOW)
conds = [
    ("1.  이기고 싶게 만드는 목표", "   경쟁, 기록 갱신, 상대 뚫기", 2.10),
    ("2.  실패해도 괜찮은 환경",   "   틀려도 감점 없음, 재시도 가능", 3.10),
    ("3.  AI를 활용할 수 있는 자유", "   방법 제한 없음, 도구 제한 없음", 4.10),
]
for main, sub, top in conds:
    tb(s, main, 0.40, top, 5.80, 0.44, size=18, color=GREEN)
    tb(s, sub,  0.40, top + 0.44, 5.80, 0.38, size=16, color=GREEN_DIM)

rect(s, 6.50, 1.50, 0.04, 5.00, GREEN_LINE)

# 오른쪽: 교과별 예시
tb(s, "교과별 적용 예시", 6.70, 1.55, 6.20, 0.44, size=19, bold=True, color=YELLOW)
examples = [
    ("▸  국어", "가장 설득력 있는 글을 AI로 써서 대결", 2.10),
    ("▸  수학", "AI로 문제를 만들고 서로 풀기", 3.10),
    ("▸  과학", "AI로 실험 설계하고 결과 예측 대결", 4.10),
]
for main, sub, top in examples:
    tb(s, main, 6.70, top, 6.20, 0.44, size=18, color=GREEN)
    tb(s, sub,  6.70, top + 0.44, 6.20, 0.38, size=16, color=GREEN_DIM)

# ════════════════════════════════════════════════════════════
# 슬라이드 27 — 정리
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "정리", pagenum="27 / 30")

bullet(s, "▸  AI를 활용하면 코딩 경험이 없어도 생각 이상의 성과를 낸다", 1.55)
bullet(s, "▸  배경 지식의 차이보다 동기의 차이가 결과를 가른다", 2.10)
bullet(s, "▸  학생 스스로 문제를 찾고 해결하게 만드는 환경이 핵심이다", 2.65)
bullet(s, "▸  교사의 역할은 가르치는 것이 아니라 하고 싶게 만드는 것이다", 3.20)

rect(s, 0.85, 4.30, 9.50, 1.60, GREEN_DARK)
rect(s, 0.85, 4.30, 0.06, 1.60, GREEN)
tb(s, "AI 시대, 교사가 해야 할 일\n학생이 해보고 싶게 만드는 것",
   1.05, 4.45, 9.10, 1.30, size=26, bold=True, color=GREEN)

# ════════════════════════════════════════════════════════════
# 슬라이드 28 — 오늘 연수 안내
# ════════════════════════════════════════════════════════════
s = new_slide()
header(s, "오늘 연수 안내", "지금까지: 1차시 사례 발표 완료", "28 / 30")

bullet(s, "앞으로", 1.55, color=YELLOW)

card(s, 0.40, 2.10, 5.80, 2.60, GREEN, "2차시", "AI와 함께 문제 해결 방법 탐색",
     ["▸  참가자 학생 입장으로 전환",
      "▸  AI 도구 활용해 스스로 과제 해결 방법 탐색"])

card(s, 6.90, 2.10, 5.80, 2.60, YELLOW, "3차시", "참가자 간 실전 시합",
     ["▸  실시간 순위 공개",
      "▸  상위 참가자 소정의 상품 증정 이벤트"])

tb(s, "지금 체험할 사이트",
   0.50, 5.10, 12.30, 0.44, size=17, color=GREEN_DIM)
tb(s, "https://macro-classroom.shk-8b6.workers.dev/",
   0.50, 5.55, 12.30, 0.44, size=19, color=GREEN)

# ════════════════════════════════════════════════════════════
# 슬라이드 29 — Q&A
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 1.50, 1.80, 10.30, 4.00, TITLE_BG)
rect(s, 4.00, 4.15, 5.30, 0.03, GREEN_DIM)
tb(s, "Q & A", 1.60, 2.20, 10.00, 1.40, size=64, bold=True, color=GREEN,
   align=PP_ALIGN.CENTER)
tb(s, "질문 받겠습니다", 1.60, 3.70, 10.00, 0.60, size=24, color=GREEN_DIM,
   align=PP_ALIGN.CENTER)
tb(s, "29 / 30", 12.20, 7.10, 1.00, 0.30, size=11, color=GRAY)

# ════════════════════════════════════════════════════════════
# 슬라이드 30 — 감사합니다
# ════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 1.50, 1.80, 10.30, 4.00, TITLE_BG)
rect(s, 4.00, 4.15, 5.30, 0.03, GREEN_DIM)
tb(s, "감사합니다", 1.60, 2.20, 10.00, 1.20, size=54, bold=True, color=GREEN,
   align=PP_ALIGN.CENTER)
tb(s, "가르치지 않아도 배운다", 1.60, 3.55, 10.00, 0.60, size=22, color=GREEN_DIM,
   align=PP_ALIGN.CENTER)
tb(s, "30 / 30", 12.20, 7.10, 1.00, 0.30, size=11, color=GRAY)

# ════════════════════════════════════════════════════════════
out = r"C:\Users\USER\Documents\GitHub\Macro_Classroom\가르치지않아도배운다_발표_v2.pptx"
prs.save(out)
print(f"완료: {len(list(prs.slides))}슬라이드")
